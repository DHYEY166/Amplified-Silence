#!/usr/bin/env python3
"""Build Amplified Silence DSCI-531 presentation from USC template."""

from __future__ import annotations
from copy import deepcopy
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import copy

ROOT = Path(__file__).resolve().parent
TEMPLATE = ROOT / "USC_PP_Template_General_National2_16x9.pptx"
FIG_DIR = ROOT / "data" / "outputs" / "figures"
OUT = ROOT / "AmplifiedSilence_Presentation.pptx"

# USC brand colors
USC_CARDINAL = RGBColor(0x99, 0x00, 0x00)
USC_GOLD = RGBColor(0xFF, 0xCC, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_GRAY = WHITE          # slides 2-7 have red bg — use white for body text
MID_GRAY = RGBColor(0xDD, 0xDD, 0xDD)  # light gray on red bg

prs = Presentation(TEMPLATE)
W = prs.slide_width   # 13.33 in
H = prs.slide_height  # 7.5 in

LAYOUT_TITLE = prs.slide_layouts[0]          # TITLE (center title + subtitle)
LAYOUT_CONTENT = prs.slide_layouts[1]        # Title and Content
LAYOUT_SECTION = prs.slide_layouts[2]        # Section Divider
LAYOUT_PHOTO = prs.slide_layouts[3]          # 2-col photo


def _tf_style(tf, size_pt: int = 20, bold: bool = False, color: RGBColor = DARK_GRAY,
              align=PP_ALIGN.LEFT) -> None:
    """Apply uniform style to every run in a text frame."""
    for para in tf.paragraphs:
        para.alignment = align
        for run in para.runs:
            run.font.size = Pt(size_pt)
            run.font.bold = bold
            run.font.color.rgb = color


def add_textbox(slide, text: str, left, top, width, height,
                size_pt=18, bold=False, color=DARK_GRAY,
                align=PP_ALIGN.LEFT, wrap=True) -> None:
    txb = slide.shapes.add_textbox(left, top, width, height)
    txb.text_frame.word_wrap = wrap
    tf = txb.text_frame
    para = tf.paragraphs[0]
    para.alignment = align
    run = para.add_run()
    run.text = text
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.color.rgb = color


def add_bullet_box(slide, bullets: list[str], left, top, width, height,
                   size_pt=18, bold_first=False, color=DARK_GRAY,
                   indent_pt=12) -> None:
    from pptx.util import Pt as PPt
    from pptx.oxml.ns import qn
    import lxml.etree as etree

    txb = slide.shapes.add_textbox(left, top, width, height)
    tf = txb.text_frame
    tf.word_wrap = True
    for idx, bullet in enumerate(bullets):
        if idx == 0:
            para = tf.paragraphs[0]
        else:
            para = tf.add_paragraph()
        para.alignment = PP_ALIGN.LEFT
        # set space before
        pPr = para._pPr
        if pPr is None:
            pPr = para._p.get_or_add_pPr()
        pPr.set('indent', str(int(-Pt(indent_pt))))
        pPr.set('marL', str(int(Pt(indent_pt))))
        run = para.add_run()
        run.text = ("• " if not bullet.startswith("  ") else "  – ") + bullet.lstrip()
        run.font.size = Pt(size_pt)
        run.font.bold = (bold_first and idx == 0)
        run.font.color.rgb = color


# ─────────────────────────────────────────────────────────────────────────────
# Remove all template placeholder slides; start fresh on a clean template copy
# (we reuse the template's master/theme/layouts but add our own slides)
# ─────────────────────────────────────────────────────────────────────────────

# Delete all 10 existing slides from the template copy
from pptx.oxml.ns import qn as _qn
R_NS = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'
sldIdLst = prs.slides._sldIdLst
for _ in range(len(prs.slides)):
    sld_id = sldIdLst[0]
    rId = sld_id.get(f'{{{R_NS}}}id')
    prs.part.drop_rel(rId)
    sldIdLst.remove(sld_id)

assert len(prs.slides) == 0


def new_slide(layout):
    return prs.slides.add_slide(layout)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ══════════════════════════════════════════════════════════════════════════════
s1 = new_slide(LAYOUT_TITLE)
ph = s1.placeholders
# idx 0 = center title, idx 1 = subtitle
for p in s1.placeholders:
    if p.placeholder_format.idx == 0:
        tf = p.text_frame
        tf.clear()
        para = tf.paragraphs[0]
        para.alignment = PP_ALIGN.LEFT
        run = para.add_run()
        run.text = "Amplified Silence"
        run.font.size = Pt(36)
        run.font.bold = True
        run.font.color.rgb = WHITE
        # subtitle line
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.LEFT
        r2 = p2.add_run()
        r2.text = "Uncovering and Mitigating Gender Bias in Music Recommendation Algorithms"
        r2.font.size = Pt(22)
        r2.font.bold = False
        r2.font.color.rgb = USC_GOLD
    elif p.placeholder_format.idx == 1:
        tf = p.text_frame
        tf.clear()
        para = tf.paragraphs[0]
        para.alignment = PP_ALIGN.LEFT
        run = para.add_run()
        run.text = "Dhyey Desai  ·  Shivam Bhosale  ·  Tejas Jonnalagadda"
        run.font.size = Pt(16)
        run.font.bold = False
        run.font.color.rgb = WHITE
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.LEFT
        r2 = p2.add_run()
        r2.text = "DSCI 531 — Fairness in AI  |  University of Southern California"
        r2.font.size = Pt(14)
        r2.font.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — PROBLEM DESCRIPTION
# ══════════════════════════════════════════════════════════════════════════════
s2 = new_slide(LAYOUT_CONTENT)
for p in s2.placeholders:
    if p.placeholder_format.idx == 0:
        tf = p.text_frame
        tf.clear()
        para = tf.paragraphs[0]
        run = para.add_run()
        run.text = "Problem Description"
        run.font.size = Pt(28)
        run.font.bold = True
        run.font.color.rgb = USC_GOLD
    elif p.placeholder_format.idx == 1:
        tf = p.text_frame
        tf.clear()
        lines = [
            ("Gender Bias in Music Recommendations", True, 14, USC_GOLD),
            ("Music streaming platforms surface artists, not just songs — who gets recommended shapes careers", False, 13, DARK_GRAY),
            ("Female and non-binary artists receive systematically less algorithmic visibility", False, 13, DARK_GRAY),
            ("", False, 1, DARK_GRAY),
            ("Why It Matters: Provider-Side Fairness", True, 14, USC_GOLD),
            ("Traditional fairness focuses on users — we focus on artists (providers)", False, 13, DARK_GRAY),
            ("Even a small recommendation bias → compounded exposure gap at scale", False, 13, DARK_GRAY),
            ("Popularity feedback loop amplifies initial skew over time", False, 13, DARK_GRAY),
            ("", False, 1, DARK_GRAY),
            ("Dataset: Last.fm 360K", True, 14, USC_GOLD),
            ("24,509 users · 10,362 artists · 1.01M play events (implicit feedback)", False, 13, DARK_GRAY),
            ("Artist gender labels fetched from MusicBrainz → 3,867 artists labeled M/F", False, 13, DARK_GRAY),
            ("Training data: 72.8% male-artist play share among labeled artists", False, 13, DARK_GRAY),
        ]
        first = True
        for (text, bold, size, color) in lines:
            if first:
                para = tf.paragraphs[0]
                first = False
            else:
                para = tf.add_paragraph()
            run = para.add_run()
            run.text = ("• " if not bold and text else "") + text
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.color.rgb = color


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — METHODOLOGY: Data + Model
# ══════════════════════════════════════════════════════════════════════════════
s3 = new_slide(LAYOUT_CONTENT)
for p in s3.placeholders:
    if p.placeholder_format.idx == 0:
        tf = p.text_frame
        tf.clear()
        para = tf.paragraphs[0]
        run = para.add_run()
        run.text = "Methodology: Data Pipeline & Model"
        run.font.size = Pt(28)
        run.font.bold = True
        run.font.color.rgb = USC_GOLD
    elif p.placeholder_format.idx == 1:
        tf = p.text_frame
        tf.clear()
        lines = [
            ("Data Pipeline", True, 14, USC_GOLD),
            ("HuggingFace Last.fm 360K → filter: ≥25 plays/user, ≥15 plays/artist", False, 13, DARK_GRAY),
            ("Per-user 80/20 train-test split; log-scaled play counts as ALS confidence weights", False, 13, DARK_GRAY),
            ("MusicBrainz API (rate-limited 1 req/s): artist type → gender label {male, female, unknown}", False, 13, DARK_GRAY),
            ("11,000 API calls → 3,867 M/F labeled items (37.3% of catalog)", False, 13, DARK_GRAY),
            ("", False, 1, DARK_GRAY),
            ("Recommendation Model: ALS (Alternating Least Squares)", True, 14, USC_GOLD),
            ("Implicit-feedback collaborative filtering via the implicit library", False, 13, DARK_GRAY),
            ("Hyper-parameters: 64 latent factors, 20 iterations, λ = 0.08", False, 13, DARK_GRAY),
            ("Candidate pool of 400 artists per user; post-processed to top-10", False, 13, DARK_GRAY),
            ("", False, 1, DARK_GRAY),
            ("Training-Data Bias Observed", True, 14, USC_GOLD),
            ("Male artists dominate play history: 72.8% of weighted interactions", False, 13, DARK_GRAY),
            ("This upstream skew feeds directly into ALS latent factors", False, 13, DARK_GRAY),
        ]
        first = True
        for (text, bold, size, color) in lines:
            if first:
                para = tf.paragraphs[0]
                first = False
            else:
                para = tf.add_paragraph()
            run = para.add_run()
            run.text = ("• " if not bold and text else "") + text
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.color.rgb = color


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — METHODOLOGY: Fairness Metrics + Mitigations
# ══════════════════════════════════════════════════════════════════════════════
s4 = new_slide(LAYOUT_CONTENT)
for p in s4.placeholders:
    if p.placeholder_format.idx == 0:
        tf = p.text_frame
        tf.clear()
        para = tf.paragraphs[0]
        run = para.add_run()
        run.text = "Methodology: Fairness Metrics & Mitigations"
        run.font.size = Pt(28)
        run.font.bold = True
        run.font.color.rgb = USC_GOLD
    elif p.placeholder_format.idx == 1:
        tf = p.text_frame
        tf.clear()
        lines = [
            ("Provider-Side Fairness Metrics (all computed over top-10 lists)", True, 14, USC_GOLD),
            ("Exposure Male Share: fraction of position-weighted exposure going to male artists", False, 13, DARK_GRAY),
            ("Avg Rank Gap (M−F): mean rank of male items minus mean rank of female items", False, 13, DARK_GRAY),
            ("Coverage Gap (M−F): distinct male artists recommended minus distinct female artists", False, 13, DARK_GRAY),
            ("", False, 1, DARK_GRAY),
            ("Post-Processing Mitigation Strategies (no model retraining)", True, 14, USC_GOLD),
            ("Mitigation A — Score Boost +12%: multiply female artist scores by 1.12 before re-ranking", False, 13, DARK_GRAY),
            ("Mitigation B — Score Boost +28%: multiply female artist scores by 1.28 before re-ranking", False, 13, DARK_GRAY),
            ("Mitigation C — Constrained Re-rank: force ≥4 female artists into each top-10 list", False, 13, DARK_GRAY),
            ("", False, 1, DARK_GRAY),
            ("Utility Metrics (user-side, for accuracy/fairness tradeoff)", True, 14, USC_GOLD),
            ("Precision@10, Recall@10, NDCG@10 — computed over held-out test interactions", False, 13, DARK_GRAY),
            ("Sensitivity sweep: boost ∈ {0, 0.05, 0.10, 0.15, 0.20, 0.30, 0.45}", False, 13, DARK_GRAY),
        ]
        first = True
        for (text, bold, size, color) in lines:
            if first:
                para = tf.paragraphs[0]
                first = False
            else:
                para = tf.add_paragraph()
            run = para.add_run()
            run.text = ("• " if not bold and text else "") + text
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.color.rgb = color


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — RESULTS: Exposure + Amplification (fig1 + fig4)
# ══════════════════════════════════════════════════════════════════════════════
s5 = new_slide(LAYOUT_CONTENT)
for p in s5.placeholders:
    if p.placeholder_format.idx == 0:
        tf = p.text_frame
        tf.clear()
        para = tf.paragraphs[0]
        run = para.add_run()
        run.text = "Results: Gender Exposure & Bias Amplification"
        run.font.size = Pt(28)
        run.font.bold = True
        run.font.color.rgb = USC_GOLD
    elif p.placeholder_format.idx == 1:
        # We'll add figures manually; clear this placeholder
        tf = p.text_frame
        tf.clear()
        para = tf.paragraphs[0]
        run = para.add_run()
        run.text = " "

# Add two figures side by side — smaller and pushed down to avoid title overlap
fig1 = str(FIG_DIR / "fig1_exposure_male_share.png")
fig4 = str(FIG_DIR / "fig4_amplification.png")
s5.shapes.add_picture(fig1, Inches(0.3), Inches(2.1), Inches(5.8), Inches(3.7))
s5.shapes.add_picture(fig4, Inches(6.5), Inches(2.1), Inches(5.6), Inches(3.7))

# Figure captions — kept short, above USC logo
add_textbox(s5,
    "Fig 1: Position-weighted male exposure share per model. Dashed = training data (72.8%); dotted = parity (50%).",
    Inches(0.3), Inches(5.9), Inches(5.8), Inches(0.55),
    size_pt=10, bold=False, color=MID_GRAY)
add_textbox(s5,
    "Fig 2: Male exposure share at each pipeline stage — shows how score boosting progressively corrects training skew.",
    Inches(6.5), Inches(5.9), Inches(5.6), Inches(0.55),
    size_pt=10, bold=False, color=MID_GRAY)

# Key callout in gold — visible on red background
add_textbox(s5,
    "Baseline amplifies training skew: 72.8% → 69.8% male exposure.  "
    "Mit B achieves near-parity at 49.2%.",
    Inches(0.3), Inches(1.45), Inches(12.7), Inches(0.55),
    size_pt=13, bold=False, color=USC_GOLD)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — RESULTS: Fairness–Utility Tradeoff (fig3 + numbers table)
# ══════════════════════════════════════════════════════════════════════════════
s6 = new_slide(LAYOUT_CONTENT)
for p in s6.placeholders:
    if p.placeholder_format.idx == 0:
        tf = p.text_frame
        tf.clear()
        para = tf.paragraphs[0]
        run = para.add_run()
        run.text = "Results: Fairness–Utility Tradeoff"
        run.font.size = Pt(28)
        run.font.bold = True
        run.font.color.rgb = USC_GOLD
    elif p.placeholder_format.idx == 1:
        tf = p.text_frame
        tf.clear()
        para = tf.paragraphs[0]
        run = para.add_run()
        run.text = " "

# Tradeoff figure (left) — smaller, pushed down
fig3 = str(FIG_DIR / "fig3_fairness_utility_tradeoff.png")
s6.shapes.add_picture(fig3, Inches(0.2), Inches(1.9), Inches(6.2), Inches(3.9))

# Figure caption — one line, fits above USC logo
add_textbox(s6,
    "Fig 3: Each point = a boost value b. X = male exposure share (← fairer); Y = NDCG@10. "
    "Curve is flat until Mit B (b=0.28), then drops steeply — marking the practical optimum.",
    Inches(0.2), Inches(5.9), Inches(6.2), Inches(0.65),
    size_pt=10, bold=False, color=MID_GRAY)

# ── Proper pptx table (right side) ──────────────────────────────────────────
from pptx.util import Pt as _Pt
from pptx.dml.color import RGBColor as _RGB
from pptx.oxml.ns import qn as _qn2

tbl_rows_data = [
    ("Model",        "P@10",   "R@10",   "NDCG@10", "Male Exp%"),
    ("Baseline",     "0.132",  "0.156",  "0.1745",  "69.8%"),
    ("Mit A (+12%)", "0.132",  "0.155",  "0.1740",  "60.3%"),
    ("Mit B (+28%)", "0.130",  "0.154",  "0.1717",  "49.2% ★"),
    ("Mit C (hard)", "0.111",  "0.131",  "0.1095",  "14.6%"),
]

tbl_left   = Inches(6.7)
tbl_top    = Inches(2.0)
tbl_width  = Inches(6.4)
tbl_height = Inches(3.5)

tbl_shape = s6.shapes.add_table(
    len(tbl_rows_data), 5, tbl_left, tbl_top, tbl_width, tbl_height
)
tbl = tbl_shape.table

# Column widths
col_widths = [Inches(1.9), Inches(0.9), Inches(0.9), Inches(1.1), Inches(1.1)]
for ci, cw in enumerate(col_widths):
    tbl.columns[ci].width = cw

HEADER_BG  = USC_CARDINAL
HEADER_FG  = USC_GOLD
ROW_BG_ODD = RGBColor(0x6B, 0x00, 0x00)   # dark cardinal
ROW_BG_EVN = RGBColor(0x55, 0x00, 0x00)   # slightly darker
STAR_BG    = RGBColor(0x44, 0x44, 0x00)   # dark gold tint for Mit B row

for ri, row_data in enumerate(tbl_rows_data):
    for ci, cell_text in enumerate(row_data):
        cell = tbl.cell(ri, ci)
        cell.text = cell_text
        tf = cell.text_frame
        tf.word_wrap = False
        para = tf.paragraphs[0]
        para.alignment = PP_ALIGN.CENTER
        run = para.runs[0] if para.runs else para.add_run()
        run.font.size = Pt(12)
        run.font.bold = (ri == 0 or ri == 3)  # header + Mit B bold
        run.font.color.rgb = USC_GOLD if ri == 0 else WHITE

        # Cell background
        from pptx.oxml import parse_xml
        from pptx.oxml.ns import nsmap
        if ri == 0:
            bg_color = HEADER_BG
        elif ri == 3:          # Mit B — highlight row
            bg_color = RGBColor(0x2A, 0x2A, 0x00)
        elif ri % 2 == 1:
            bg_color = ROW_BG_ODD
        else:
            bg_color = ROW_BG_EVN

        tc = cell._tc
        tcPr = tc.get_or_add_tcPr()
        solidFill_xml = (
            f'<a:solidFill xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
            f'<a:srgbClr val="{str(bg_color)}"/></a:solidFill>'
        )
        solidFill = parse_xml(solidFill_xml)
        existing = tcPr.find(
            '{http://schemas.openxmlformats.org/drawingml/2006/main}solidFill'
        )
        if existing is not None:
            tcPr.remove(existing)
        tcPr.insert(0, solidFill)

# Callout below table
add_textbox(s6,
    "★  Mit B: near-parity (49.2%) with only 1.7% NDCG drop — practical sweet spot",
    tbl_left, Inches(5.65), tbl_width, Inches(0.6),
    size_pt=13, bold=True, color=USC_GOLD)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — TAKEAWAYS & LIMITATIONS
# ══════════════════════════════════════════════════════════════════════════════
s7 = new_slide(LAYOUT_CONTENT)
# Only write the title placeholder; body goes into a manual textbox to avoid top-gap
for p in s7.placeholders:
    if p.placeholder_format.idx == 0:
        tf = p.text_frame
        tf.clear()
        para = tf.paragraphs[0]
        run = para.add_run()
        run.text = "Key Takeaways & Limitations"
        run.font.size = Pt(28)
        run.font.bold = True
        run.font.color.rgb = USC_GOLD
    elif p.placeholder_format.idx == 1:
        p.text_frame.clear()

# Manual textbox starting right below the title
s7_body = s7.shapes.add_textbox(Inches(0.5), Inches(1.95), Inches(12.5), Inches(5.2))
s7_body.text_frame.word_wrap = True
lines = [
    ("Key Findings", True, 14, USC_GOLD),
    ("ALS amplifies training bias: 72.8% → 69.8% male exposure at baseline", False, 12, WHITE),
    ("Score boosting is effective and cheap: no retraining required", False, 12, WHITE),
    ("Mitigation B (+28%) is the practical sweet spot: near-parity with <2% NDCG loss", False, 12, WHITE),
    ("Hard constraint (Mit C) achieves strongest fairness but collapses NDCG by 37%", False, 12, WHITE),
    ("", False, 9, WHITE),
    ("Limitations", True, 14, USC_GOLD),
    ("Binary gender framing — non-binary, agender identities not represented", False, 12, WHITE),
    ("MusicBrainz labels only 37.3% of artists — fairness metrics exclude 62.7% of catalog", False, 12, WHITE),
    ("Post-processing only — upstream bias in training data remains unaddressed", False, 12, WHITE),
    ("No user-side fairness analysis; user demographics not available", False, 12, WHITE),
    ("", False, 9, WHITE),
    ("Future Work", True, 14, USC_GOLD),
    ("In-processing fairness constraints during ALS optimization", False, 12, WHITE),
    ("Beyond binary: multi-attribute fairness (genre, geography, label size)", False, 12, WHITE),
    ("Longitudinal study: does boosting sustain coverage long-term?", False, 12, WHITE),
]
first = True
for (text, bold, size, color) in lines:
    if first:
        para = s7_body.text_frame.paragraphs[0]
        first = False
    else:
        para = s7_body.text_frame.add_paragraph()
    run = para.add_run()
    run.text = ("• " if not bold and text else "") + text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    # add breathing room between bullet points (not spacer rows)
    if not bold and text:
        para.space_before = Pt(4)


prs.save(OUT)
print(f"Saved: {OUT}")
